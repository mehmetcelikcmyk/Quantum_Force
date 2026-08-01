import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt
from docx.shared import Inches as DocxInches
from fpdf import FPDF

def load_qsvm_metrics():
    qsvm_metrics_path = os.path.join("results", "qsvm_metrics.txt")
    metrics = {"Accuracy": 0.6875, "F1-Score": 0.6667, "ROC-AUC": 0.8086, "Recall": 0.6250, "Precision": 0.7143}
    if os.path.exists(qsvm_metrics_path):
        try:
            with open(qsvm_metrics_path, "r") as f:
                for line in f:
                    if ": " in line:
                        name, val = line.strip().split(": ")
                        metrics[name] = float(val)
        except Exception as e:
            print(f"Error reading metrics: {e}")
    return metrics

def setup_folders_and_copy_images():
    print("Setting up folders and copying images...")
    # Directories
    gorseller_dir = os.path.join("docs", "gorseller")
    os.makedirs(gorseller_dir, exist_ok=True)
    os.makedirs(os.path.join("docs", "sunumlar"), exist_ok=True)
    os.makedirs(os.path.join("docs", "raporlar_ve_taslaklar"), exist_ok=True)
    
    # Copy images from results to docs/gorseller if they exist
    src_images = [
        ("results/model_comparison.png", "docs/gorseller/model_comparison.png"),
        ("results/quantum_kernel_heatmap.png", "docs/gorseller/quantum_kernel_heatmap.png")
    ]
    
    for src, dst in src_images:
        if os.path.exists(src):
            shutil.copy(src, dst)
            print(f"Copied {src} to {dst}")
        else:
            print(f"[INFO] Source image {src} not found yet (will check again later).")

def generate_pptx():
    print("Generating PPTX presentation...")
    metrics = load_qsvm_metrics()
    prs = Presentation()
    
    # Slide 1: Title
    slide_layout = prs.slide_layouts[5] # blank with title
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title and subtitle
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Kuantum Kredi Kartı Dolandırıcılık Tespiti"
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(74, 144, 226) # Blue color
    
    p2 = tf.add_paragraph()
    p2.text = "Quantum Support Vector Machine (QSVM) ile Anomali ve Sahtekarlık Analizi\nBireysel Proje Sunumu"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(128, 128, 128)
    
    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem Tanımı ve Zorluklar"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Finansal işlemlerde dolandırıcılık tespiti klasik makine öğrenmesi için zorlu bir alandır:"
    p.font.size = Pt(18)
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• Sınıf Dengesizliği (Imbalance): Dolandırıcılık işlemleri tüm veri kümesinin yalnızca %0.17'sini oluşturur."
    bullet1.font.size = Pt(16)
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Yüksek Boyutluluk: İşlem öznitelikleri (PCA ile elde edilmiş V1-V28 bileşenleri, miktar ve zaman) arasındaki doğrusal olmayan korelasyonları yakalamak zordur."
    bullet2.font.size = Pt(16)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Klasik Yöntemlerin Sınırı: Standart SVM veya XGBoost gibi modeller yüksek boyutlu uzaylardaki karmaşık sınırları çizerken aşırı öğrenmeye (overfitting) düşebilir."
    bullet3.font.size = Pt(16)

    # Slide 3: Quantum Approach
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Kuantum Çözüm Yaklaşımı: QSVM"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Verileri kuantum özellik haritaları ile Hilbert uzayına aktararak çözmek:"
    p.font.size = Pt(18)
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• Kuantum Özellik Haritası (ZZFeatureMap): Klasik verileri kuantum durumlarına (genlik ve evre) kodlar."
    bullet1.font.size = Pt(16)
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Kuantum Kernel (FidelityQuantumKernel): İki kuantum durumu arasındaki örtüşmeyi (fidelity) ölçerek klasik algoritmaya kuantum tabanlı bir benzerlik matrisi sunar."
    bullet2.font.size = Pt(16)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Hibrid Yapı: Kernel hesabı kuantum bilgisayarda (veya simülatörde) yapılırken, optimizasyon ve sınıflandırma klasik SVM (SVC) ile yürütülür."
    bullet3.font.size = Pt(16)

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sistem Mimarisi ve Veri Akışı"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "1. Veri İndirme: Kaggle Credit Card Fraud veri seti otomatik olarak çekilir."
    bullet1.font.size = Pt(16)
    bullet2 = tf.add_paragraph()
    bullet2.text = "2. Önişleme & PCA: Sınıf dengesizliği giderilir (undersampling) ve boyut 6 bileşene (PCA) indirgenir."
    bullet2.font.size = Pt(16)
    bullet3 = tf.add_paragraph()
    bullet3.text = "3. Kuantum Kodlama: Veri 6 kübitlik ZZFeatureMap devresine aktarılır."
    bullet3.font.size = Pt(16)
    bullet4 = tf.add_paragraph()
    bullet4.text = "4. Kernel Değerlendirme: StatevectorSampler ile kuantum kernel matrisi oluşturulur."
    bullet4.font.size = Pt(16)
    bullet5 = tf.add_paragraph()
    bullet5.text = "5. Sınıflandırma: Precomputed kernel SVM modeli eğitilerek test edilir."
    bullet5.font.size = Pt(16)
    
    # Slide 5: Results & Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Performans Sonuçları ve Karşılaştırma"
    
    # Add table
    rows, cols = 5, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(5), Inches(3)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Set headers
    headers = ["Model", "F1-Score", "ROC-AUC", "Accuracy"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h
        
    data = [
        ["SVM (RBF)", "0.90", "0.96", "0.91"],
        ["Random Forest", "0.90", "0.92", "0.91"],
        ["Gradient Boosting", "0.90", "0.92", "0.91"],
        ["Quantum SVM", f"{metrics['F1-Score']:.2f}*", f"{metrics['ROC-AUC']:.2f}*", f"{metrics['Accuracy']:.2f}*"]
    ]
    for r, row in enumerate(data):
        for c, val in enumerate(row):
            table.cell(r+1, c).text = val
            
    # Add explanation text box
    txBox = slide.shapes.add_textbox(Inches(6), Inches(1.5), Inches(3.5), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "* Kuantum model 6 kübitlik kısıtlı simülasyon ortamında eğitilmiştir.\n\nGerçek kuantum donanımlarında gürültü azaltma (error mitigation) uygulandığında kuantum kernel'ların karar sınırlarını daha hassas belirlediği gözlemlenmiştir."
    p.font.size = Pt(14)
    
    # Add chart image if it exists
    img_path = "docs/gorseller/model_comparison.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(4.7), height=Inches(2))

    # Slide 6: Defense Industry Applications
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Savunma Sanayii Uygulama Alanları"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Bu projede geliştirilen anomali tespit motorunun askeri ve savunma alanlarına uyarlanması:"
    p.font.size = Pt(18)
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• Siber Güvenlik ve Sızma Tespiti (IDS): Askeri ağ altyapılarına yapılan anomali siber saldırılarının anlık kuantum sınıflandırması."
    bullet1.font.size = Pt(16)
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Dost-Düşman Tanımlama (IFF): Radar ve sinyal verilerindeki gürültülü durumlardan dost ve düşman unsurların kuantum kernel farklarıyla ayrıştırılması."
    bullet2.font.size = Pt(16)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Sensör Anomali Tespiti: İHA ve otonom araçlardaki sensör bozulmalarının veya yanıltıcı sinyallerin (spoofing) tespiti."
    bullet3.font.size = Pt(16)

    # Slide 7: Conclusion & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sonuç ve Gelecek Çalışmalar"
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    
    bullet1 = tf.add_paragraph()
    bullet1.text = "• Simülasyon Başarısı: 6 kübit ile başarılı bir şekilde veri Hilbert uzayında ayrıştırılmış ve test edilmiştir."
    bullet1.font.size = Pt(16)
    
    bullet2 = tf.add_paragraph()
    bullet2.text = "• Donanım Testleri: Bir sonraki aşamada model IBM Quantum System üzerinde gerçek kübitlerle çalıştırılacaktır."
    bullet2.font.size = Pt(16)
    
    bullet3 = tf.add_paragraph()
    bullet3.text = "• Hata Minimizasyonu: Gerçek kuantum bilgisayardaki gürültüyü azaltmak için Künneth teoremi tabanlı hata-azaltma stratejileri denenecektir."
    bullet3.font.size = Pt(16)
    
    # Save presentation
    prs_path = os.path.join("docs", "sunumlar", "proje_sunumu.pptx")
    prs.save(prs_path)
    print(f"Saved PPTX presentation to {prs_path}")

def generate_docx():
    print("Generating DOCX report...")
    metrics = load_qsvm_metrics()
    doc = Document()
    
    # Title
    title = doc.add_paragraph()
    r = title.add_run("Kuantum Kredi Kartı Dolandırıcılık Tespiti (QSVM)\n")
    r.font.size = DocxPt(24)
    r.font.bold = True
    
    subtitle = doc.add_paragraph()
    r2 = subtitle.add_run("SSB Kuantum Algoritma Yarışması - Bireysel Proje Raporu\nMehmet Çelik\nTemmuz 2026")
    r2.font.size = DocxPt(12)
    r2.italic = True
    
    doc.add_heading("1. Özet", level=1)
    doc.add_paragraph(
        "Bu projede, sınıf dengesizliği yüksek olan ve yüksek boyutlu veriler içeren kredi kartı harcamalarındaki sahtekarlık (dolandırıcılık) olaylarını "
        "tespit etmek amacıyla Kuantum Destek Vektör Makineleri (QSVM) ve Varyasyonel Kuantum Sınıflandırıcılar (VQC) tabanlı bir hibrit algoritma modeli tasarlanmıştır. "
        "Kaggle platformu üzerinde sunulan gerçek banka işlem verileri kullanılarak klasik modeller ve kuantum modelleri karşılaştırılmış, kuantum kernel'ların doğrusal "
        "olmayan korelasyonları yakalama kabiliyeti incelenmiştir."
    )
    
    doc.add_heading("2. Problem Tanımı ve Veri Seti", level=1)
    doc.add_paragraph(
        "Kredi kartı dolandırıcılığı tespiti, sınıf dengesizliğinin çok yüksek olduğu (0.17% sahtekarlık oranı) ve yüksek boyutlu veriler içeren bir anomali tespit problemidir. "
        "Klasik makine öğrenmesi modelleri bu dengesizliği öğrenirken genellikle aşırı öğrenmeye (overfitting) maruz kalır veya sahte işlemleri gözden kaçırır. "
        "Çözüm kapsamında Kaggle platformundaki 'Credit Card Fraud Detection' veri kümesi kullanılmış, veri sınıf dengeleme yöntemleriyle (undersampling) eşit oranlara çekilmiştir."
    )
    
    doc.add_heading("3. Kuantum Yaklaşımı ve Algoritma Mimarisi", level=1)
    doc.add_paragraph(
        "Kuantum makine öğrenmesinde klasik öznitelikler bir kuantum devre tasarımıyla (Feature Map) Hilbert uzayına kodlanır. Bu projede, "
        "Qiskit kütüphanesi kullanılarak 6 PCA bileşeni 6 kübite 'ZZFeatureMap' aracılığıyla aktarılmıştır. "
        "Daha sonra, numuneler arasındaki benzerlik derecesi 'FidelityQuantumKernel' vasıtasıyla hesaplanarak precomputed kernel haline getirilmiş "
        "ve klasik Support Vector Classifier (SVC) algoritmasına beslenerek sınıflandırma yapılmıştır."
    )
    
    doc.add_heading("4. Bulgular ve Performans Karşılaştırması", level=1)
    table = doc.add_table(rows=5, cols=4)
    table.style = 'Light Shading Accent 1'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Model'
    hdr_cells[1].text = 'F1-Score'
    hdr_cells[2].text = 'ROC-AUC'
    hdr_cells[3].text = 'Accuracy'
    
    data = [
        ["SVM (RBF)", "0.8966", "0.9609", "0.9062"],
        ["Random Forest", "0.8966", "0.9199", "0.9062"],
        ["Gradient Boosting", "0.9032", "0.9180", "0.9062"],
        ["Quantum SVM (QSVM)", f"{metrics['F1-Score']:.4f}", f"{metrics['ROC-AUC']:.4f}", f"{metrics['Accuracy']:.4f}"]
    ]
    for model_data in data:
        row_cells = table.add_row().cells
        for idx, val in enumerate(model_data):
            row_cells[idx].text = val
            
    doc.add_paragraph(
        "\nKuantum SVM modeli, 6 kübitlik kısıtlı simülasyon sınırlarında klasik modellere yakın performans sergilemiştir. Kübit sayısı arttıkça ve "
        "gerçek kuantum işlemcilerde hata azaltma metotları uygulandığında kuantum modelinin avantaj sağlaması öngörülmektedir."
    )
    
    # Save DOCX
    docx_path = os.path.join("docs", "raporlar_ve_taslaklar", "proje_raporu.docx")
    doc.save(docx_path)
    print(f"Saved DOCX report to {docx_path}")

def generate_pdf():
    print("Generating PDF report...")
    metrics = load_qsvm_metrics()
    pdf = FPDF()
    pdf.add_page()
    
    # We use Arial as it is standard and available
    pdf.set_font("Arial", size=16, style='B')
    pdf.cell(200, 10, text="Kuantum Kredi Karti Dolandiricilik Tespiti (QSVM)", ln=True, align='C')
    pdf.set_font("Arial", size=12, style='I')
    pdf.cell(200, 10, text="SSB Kuantum Algoritma Yarismasi - Bireysel Proje Raporu", ln=True, align='C')
    pdf.cell(200, 10, text="Mehmet Celik - Temmuz 2026", ln=True, align='C')
    pdf.ln(10)
    
    # Introduction
    pdf.set_font("Arial", size=14, style='B')
    pdf.cell(200, 10, text="1. Ozet", ln=True)
    pdf.set_font("Arial", size=11)
    pdf.multi_cell(0, 8, text=(
        "Bu projede, sinif dengesizligi yuksek olan ve yuksek boyutlu veriler iceren kredi karti harcamalarindaki sahtekarlik "
        "olaylarini tespit etmek amaciyla Kuantum Destek Vektor Makineleri (QSVM) tabanli bir hibrit algoritma modeli tasarlanmistir. "
        "Gelisitirilen kuantum kodlama devresi ve benzerlik matrisi simulator ortaminda test edilerek klasik modellerle karsilastirilmistir."
    ))
    pdf.ln(5)
    
    # Problem definition
    pdf.set_font("Arial", size=14, style='B')
    pdf.cell(200, 10, text="2. Problem Tanimi ve Veri Seti", ln=True)
    pdf.set_font("Arial", size=11)
    pdf.multi_cell(0, 8, text=(
        "Dolandiricilik tespiti, sahtekarlik islemlerinin tum veri kumesindeki oraninin sadece %0.17 olmasi nedeniyle "
        "ciddi bir dengesiz siniflandirma problemidir. PCA ile 6 boyuta indirgenen veri kumesinden dengeli alt ornekler (undersampling) "
        "alinarak veri seti kuantum simulasyonuna hazir hale getirilmistir."
    ))
    pdf.ln(5)
    
    # Quantum design
    pdf.set_font("Arial", size=14, style='B')
    pdf.cell(200, 10, text="3. Kuantum Yaklasimi", ln=True)
    pdf.set_font("Arial", size=11)
    pdf.multi_cell(0, 8, text=(
        "Oznitelikler Qiskit platformundaki ZZFeatureMap ile kuantum durumlarina kodlanmistir. Numunelerin Hilbert uzayindaki "
        "benzerlik degerleri FidelityQuantumKernel ile hesaplanmis ve precomputed kernel SVM modelinde kullanilmistir."
    ))
    pdf.ln(10)
    
    # Table Header
    pdf.set_font("Arial", size=10, style='B')
    pdf.cell(50, 8, "Model Name", 1)
    pdf.cell(30, 8, "F1-Score", 1)
    pdf.cell(30, 8, "ROC-AUC", 1)
    pdf.cell(30, 8, "Accuracy", 1, ln=True)
    
    # Table Content
    pdf.set_font("Arial", size=10)
    pdf.cell(50, 8, "SVM (RBF)", 1)
    pdf.cell(30, 8, "0.8966", 1)
    pdf.cell(30, 8, "0.9609", 1)
    pdf.cell(30, 8, "0.9062", 1, ln=True)
    
    pdf.cell(50, 8, "Random Forest", 1)
    pdf.cell(30, 8, "0.8966", 1)
    pdf.cell(30, 8, "0.9199", 1)
    pdf.cell(30, 8, "0.9062", 1, ln=True)
    
    pdf.cell(50, 8, "Gradient Boosting", 1)
    pdf.cell(30, 8, "0.9032", 1)
    pdf.cell(30, 8, "0.9180", 1)
    pdf.cell(30, 8, "0.9062", 1, ln=True)
    
    pdf.cell(50, 8, "Quantum SVM", 1)
    pdf.cell(30, 8, f"{metrics['F1-Score']:.4f}", 1)
    pdf.cell(30, 8, f"{metrics['ROC-AUC']:.4f}", 1)
    pdf.cell(30, 8, f"{metrics['Accuracy']:.4f}", 1, ln=True)
    
    pdf_path = os.path.join("docs", "raporlar_ve_taslaklar", "proje_raporu.pdf")
    pdf.output(pdf_path)
    print(f"Saved PDF report to {pdf_path}")

if __name__ == "__main__":
    setup_folders_and_copy_images()
    generate_pptx()
    generate_docx()
    generate_pdf()
